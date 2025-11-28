"""
File Handler Module
Extracts text from multiple file formats (PDF, DOCX, PPTX, TXT)
"""

import streamlit as st
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io


def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    try:
        pdf_reader = PdfReader(file)
        text = ""
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            text += f"\n--- Page {page_num + 1} ---\n{page_text}"
        return text
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")
        return None


def extract_text_from_docx(file):
    """Extract text from Word (.docx) file"""
    try:
        doc = Document(file)
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():  # Only add non-empty paragraphs
                text += para.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading DOCX: {str(e)}")
        return None


def extract_text_from_pptx(file):
    """Extract text from PowerPoint (.pptx) file"""
    try:
        prs = Presentation(file)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading PPTX: {str(e)}")
        return None


def extract_text_from_txt(file):
    """Extract text from plain text file"""
    try:
        return file.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error reading TXT: {str(e)}")
        return None


def extract_text_from_file(uploaded_file):
    """
    Main function: Detects file type and extracts text appropriately
    
    Args:
        uploaded_file: Streamlit UploadedFile object
    
    Returns:
        Extracted text (string) or None if error
    """
    
    if uploaded_file is None:
        return None
    
    file_extension = uploaded_file.name.split(".")[-1].lower()
    
    st.info(f"📄 Processing: {uploaded_file.name} ({file_extension.upper()})")
    
    if file_extension == "pdf":
        return extract_text_from_pdf(uploaded_file)
    
    elif file_extension == "docx":
        return extract_text_from_docx(uploaded_file)
    
    elif file_extension == "pptx":
        return extract_text_from_pptx(uploaded_file)
    
    elif file_extension == "txt":
        return extract_text_from_txt(uploaded_file)
    
    else:
        st.error(f"❌ File type '.{file_extension}' not supported. Use: PDF, DOCX, PPTX, or TXT")
        return None
